#!/usr/bin/env python3
"""2차예선 기술개발보고서 초안 PPT 생성기.

`docs/qualifier2_report_outline.md` 의 슬라이드 골격을 그대로 .pptx 로 뽑는다.
디자인은 의도적으로 최소 — 편집 출발점이 목적이고, 꾸미기는 PowerPoint 에서 한다.

사용:
    python tools/report/build_qualifier2_deck.py [출력경로]

기본 출력: docs/qualifier2_report_draft.pptx
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

FONT = "맑은 고딕"

C_TITLE = RGBColor(0x1A, 0x1A, 0x1A)
C_SUB = RGBColor(0x66, 0x66, 0x66)
C_BODY = RGBColor(0x22, 0x22, 0x22)
C_TODO = RGBColor(0xC0, 0x39, 0x2B)
C_KEY = RGBColor(0x1F, 0x4E, 0x79)
C_RULE = RGBColor(0x1F, 0x4E, 0x79)
C_TAG = RGBColor(0x88, 0x88, 0x88)

# (본문, 들여쓰기 레벨, 스타일)  스타일: None | "todo" | "key"
Body = list[tuple[str, int, str | None]]


def _tb(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def _run(para, text, size, color, bold=False):
    r = para.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return r


CONTENT_BOTTOM = 6.62  # 본문·표가 넘어서면 안 되는 하단 경계 (note/페이지번호 자리)
CHARS_PER_LINE = {0: 82, 1: 88, 2: 96}  # 12.2in 폭 기준 한 줄에 들어가는 대략적 글자수


def _body_metrics(body: Body, scale: float):
    """본문 높이를 실제 줄 수(줄바꿈 포함)로 추정한다. (높이 in, 폰트크기 dict)"""
    sizes = {0: 16 * scale, 1: 14 * scale, 2: 12 * scale}
    line_h = {0: 0.285, 1: 0.250, 2: 0.220}
    total = 0.0
    for text, level, _ in body:
        wrapped = max(1, -(-len(text) // int(CHARS_PER_LINE[level] / scale)))
        total += line_h[level] * scale * wrapped + 0.055
    return total, sizes


def add_slide(prs, tag, title, subtitle, body: Body, table=None, note=None, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 상단 태그 (섹션 표시)
    tf = _tb(slide, Inches(0.55), Inches(0.28), Inches(9.0), Inches(0.3))
    _run(tf.paragraphs[0], tag, 12, C_TAG)

    # 제목
    tf = _tb(slide, Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.75))
    _run(tf.paragraphs[0], title, 30, C_TITLE, bold=True)

    # 제목 밑줄
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.32), Inches(12.2), Emu(11430))
    line.fill.solid()
    line.fill.fore_color.rgb = C_RULE
    line.line.fill.background()
    line.shadow.inherit = False

    top = 1.45
    if subtitle:
        tf = _tb(slide, Inches(0.55), Inches(top), Inches(12.2), Inches(0.35))
        _run(tf.paragraphs[0], subtitle, 15, C_SUB, bold=True)
        top = 1.92

    # 표 높이를 먼저 확보한다 (표가 밀려나 슬라이드를 넘는 것을 막는다)
    n_cols = len(table[0]) if table else 0
    tbl_font = 10 if n_cols <= 5 else (9 if n_cols <= 7 else 8)
    row_h = 0.30 if n_cols <= 5 else 0.40  # 열이 많으면 줄바꿈으로 행이 커진다
    tbl_h = row_h * len(table) if table else 0.0
    avail = CONTENT_BOTTOM - top - (tbl_h + 0.14 if table else 0.0)

    # 본문이 남은 공간에 들어갈 때까지 폰트를 줄인다
    scale = 1.0
    if body:
        while scale > 0.68:
            need, sizes = _body_metrics(body, scale)
            if need <= avail:
                break
            scale -= 0.04
        need, sizes = _body_metrics(body, scale)

        tf = _tb(slide, Inches(0.55), Inches(top), Inches(12.2), Inches(max(need, 0.3)))
        for i, (text, level, style) in enumerate(body):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.level = level
            p.space_after = Pt(4)
            color = {"todo": C_TODO, "key": C_KEY, None: C_BODY}[style]
            bullet = {0: "■ ", 1: "– ", 2: "· "}[level] if text else ""
            _run(p, bullet + text, sizes[level], color, bold=(style == "key"))
        top += need + 0.14

    # 표
    if table:
        rows = len(table)
        gt = slide.shapes.add_table(
            rows, n_cols, Inches(0.55), Inches(top), Inches(12.2), Inches(tbl_h)
        ).table
        if col_widths:
            unit = 12.2 / sum(col_widths)
            for c, w in enumerate(col_widths):
                gt.columns[c].width = Inches(w * unit)
        for r, row in enumerate(table):
            for c, cell_text in enumerate(row):
                cell = gt.cell(r, c)
                cell.text = ""
                cell.margin_top = Pt(1)
                cell.margin_bottom = Pt(1)
                p = cell.text_frame.paragraphs[0]
                is_todo = "[빈칸]" in cell_text or "[팀장" in cell_text
                _run(p, cell_text, tbl_font, C_TODO if is_todo else C_BODY, bold=(r == 0))
                cell.text_frame.word_wrap = True

    # 하단 메모
    if note:
        tf = _tb(slide, Inches(0.55), Inches(6.78), Inches(11.2), Inches(0.4))
        _run(tf.paragraphs[0], note, 11, C_SUB)

    # 페이지 번호 (양식 필수) — 표지·목차 제외하고 1번부터
    idx = len(prs.slides) - 2
    if idx >= 1:
        tf = _tb(slide, Inches(12.35), Inches(6.98), Inches(0.75), Inches(0.3))
        p = tf.paragraphs[0]
        p.alignment = 2
        _run(p, str(idx), 11, C_SUB)

    return slide


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tf = _tb(slide, Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.0))
    _run(tf.paragraphs[0], "제24회 한국로봇항공기경연대회", 20, C_SUB)

    tf = _tb(slide, Inches(1.0), Inches(2.8), Inches(11.3), Inches(1.4))
    _run(tf.paragraphs[0], "2차 예선 기술개발보고서", 48, C_TITLE, bold=True)

    tf = _tb(slide, Inches(1.0), Inches(4.3), Inches(11.3), Inches(1.6))
    for i, line in enumerate(
        [
            "정규부문   |   참가팀명 : 수리독수리 (연세대학교)",
            "지도교수 : 유정훈 교수 (연세대학교 기계공학부)",
            "발표자 : [빈칸]        2026. 07. 30.",
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        _run(p, line, 16, C_TODO if "[빈칸]" in line else C_BODY)
    return slide


def add_toc(prs):
    body: Body = [
        ("1. 참가팀 현황", 0, "key"),
        ("참가팀 구성 / 추진체계", 1, None),
        ("2. 전체 시스템 개발 상황", 0, "key"),
        ("전체 시스템 구성 / 설계 방법론 / 주익 최적해 / 안정성 설계 /", 1, None),
        ("구조·진동 해석 / 제작·실물 검증 / 자체 제작 시험기체군 /", 1, None),
        ("미완성 개발·시험 계획 / 안전장치", 1, None),
        ("3. 자동비행 시스템 구현 상황", 0, "key"),
        ("유도·제어·항법 아키텍처 / 경로생성 / 경로추종·상태머신 /", 1, None),
        ("검증 체계 / 지상·비행 시험", 1, None),
        ("4. 임무수행 구현기술", 0, "key"),
        ("임무 시나리오 / 임무장치(아기수리) / 객체인식·정밀착륙 / 보완 계획", 1, None),
        ("5. 기타", 0, "key"),
        ("비행시험 데이터 & 영상", 1, None),
    ]
    return add_slide(prs, "", "목  차", "", body)


SEC1 = "1. 참가팀 현황"
SEC2 = "2. 전체 시스템 개발 상황"
SEC3 = "3. 자동비행 시스템 구현 상황"
SEC4 = "4. 임무수행 구현기술"
SEC5 = "5. 기타"


def build(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover(prs)
    add_toc(prs)

    # ---------------- 1. 참가팀 현황 ----------------
    add_slide(
        prs,
        SEC1,
        "참가팀 현황",
        "연세대학교 수리독수리",
        [
            ("참가팀명(소속) : 수리독수리 (연세대학교)", 0, None),
            ("지도교수 : 유정훈 교수 (연세대학교 기계공학부)", 0, None),
            ("참가팀 구성 : 김승원 팀장 외 [빈칸]명", 0, "todo"),
        ],
        table=[
            ["No", "직위", "성명", "주민번호 앞7자리", "학년", "소속", "학과", "E-mail", "담당역할"],
            ["1", "팀장", "김승원", "050106-3", "2", "연세대", "기계공학부", "suridogsuri@gmail.com", "기체설계(HW)·시스템통합(PM)"],
            ["2", "팀원", "김태훈", "031226-3", "2", "전남대", "물리학과", "potato0303kth@jnu.ac.kr", "임무제어(SW)·지상통제(GCS)"],
            ["3", "팀원", "[빈칸]", "[빈칸]", "", "", "", "", "조종사"],
        ],
        col_widths=[0.4, 0.7, 0.8, 1.4, 0.5, 0.8, 1.3, 2.2, 2.4],
        note="※ 팀원 변경은 2차 예선 제출자료 이후 불가 — 조종사 지정 여부를 제출 전 확정할 것",
    )

    add_slide(
        prs,
        SEC1,
        "추진체계",
        "2인 체제 · 전 영역 자체 수행",
        [
            ("외주 없음 — 기체 · 소프트웨어 · 시험장비까지 전부 자체 제작", 0, "key"),
            ("2인 팀이므로 1인이 복수 기능을 겸한다 (아래 분장표)", 0, None),
        ],
        table=[
            ["기능", "담당", "역할 및 담당업무"],
            ["기체설계 (HW)", "김승원", "공력·구조·진동 해석, 형상 최적 설계, 기체 제작 및 정비"],
            ["시스템통합 (PM)", "김승원", "대회 규정 분석, 개발 일정 관리, HW-SW 통합 조율"],
            ["임무제어 (SW)", "김태훈", "자율비행 소프트웨어, 객체인식, FC 파라미터 튜닝"],
            ["지상통제 (GCS)", "김태훈", "비행 경로 설정, 텔레메트리 모니터링, 비행 로그 분석"],
            ["조종사", "[빈칸]", "수동 이착륙, 비상 상황 시 수동 전환 및 안전 확보"],
        ],
    )

    # ---------------- 2. 전체 시스템 개발 상황 ----------------
    add_slide(
        prs,
        SEC2 + " > 2-1",
        "전체 시스템 구성",
        "1차 계획서의 설계안 → 2차 시점의 구현체",
        [
            ("[그림] 시스템 블록도 : 기체(수리1호) - FC(Pixhawk 6C) - CC(RPi5) - 센서(카메라·라이다·Dual GPS) - 지상국", 0, "todo"),
        ],
        table=[
            ["구성 요소", "1차 예선 (계획)", "2차 예선 (현재)"],
            ["기체 수리1호", "설계 완료", "제작 완료 · 호버링 시험 완료"],
            ["비행제어 소프트웨어", "알고리즘 선정", "실기체 배포 · 자율비행 실증 (59회분)"],
            ["객체인식", "알고리즘 설계", "온보드 실동작 · 실기체 종단간 검증"],
            ["정밀착륙 연동", "개념 설계", "인터페이스 실증 완료, 폐루프 연결 진행 중"],
            ["구조장치 아기수리", "설계 완료", "[팀장 확인]"],
        ],
    )

    add_slide(
        prs,
        SEC2 + " > 2-1",
        "설계 방법론 : 해석 기반 최적해 탐색",
        "감이 아니라 수치로 형상을 골랐다",
        [
            ("자체 개발 스프레드시트 계산기 3종을 축으로 설계 전 과정을 정량화", 0, "key"),
            ("[프로펠러] 유동속도에 따른 추력·효율 저하 정량화. Tarot 4114(320kv) / EM3115(900kv)", 1, None),
            ("공식 데이터로 정지추력·파워계수 보정 + Apps Script 가변 RPM 자동 산출", 2, None),
            ("→ 실제 추력·소비전력을 10% 오차 내로 예측", 2, "key"),
            ("[계산기] 주익 제원 입력 시 12개 레이놀즈 수 × 속도별 공력 자동 연산", 1, None),
            ("공력 371 · 푸셔 성능 144 · 중량 60개 데이터를 유기적으로 종합", 2, None),
            ("[에너지] 주익 크기 변화에 따른 중량 변동을 반영한 임무 소비 전력량 산출", 1, None),
            ("외부 해석 도구 연동 : XFLR5 (공력, 정적·동적 안정성) + ANSYS (1-way FSI, fluent, modal, harmonic)", 0, None),
            ("[그림] 전체 해석 프로세스 다이어그램 (1차 보고서 3-1)", 0, "todo"),
        ],
        note="핵심 메시지 : 설계 → 해석 → 재설계를 사람 손이 아니라 자동화된 계산 파이프라인으로 반복했다",
    )

    add_slide(
        prs,
        SEC2 + " > 2-1",
        "주익 · 에어포일 최적해 도출",
        "60개 형상 조합 전수 평가",
        [
            ("Apps Script 자동화로 루트시위 6종 × 윙스팬 10종 = 60개 조합 공력 데이터 자동 연산", 0, "key"),
            ("평가지표 : 최대 전비 순항속도 / 실속속도 / 임무별 소비 전력량", 0, None),
            ("Min-Max 정규화 + 대회 배점표 가중치 합산 → 100점 만점 환산", 1, None),
            ("제약조건을 명시적으로 부여", 0, None),
            ("배터리 70% 방전한계(170Wh) 초과 형상 배제", 1, None),
            ("가로세로비(AR) 10 이상은 '카본파이프 + ASA AERO' 플랫폼의 구조 강성 한계로 보류", 1, None),
            ("에어포일 최종 순위 : NACA2412 > SD7037 > RG15 > DAE12 > MH32 > NACA4412 > SG6043", 0, "key"),
            ("설계 실효성 보정 : 이론적 최대 양항비가 아니라, 푸셔 과부하·추진효율 저하를 반영한", 0, None),
            ("'최대 전비 속도'(km/kWh)를 기준으로 삼아 실제 최대 항속거리를 보장", 1, None),
            ("검증 결과 : 총 소비 전력량이 6S 10000mAh의 80% 가용량(192Wh)을 안정적으로 하회", 0, "key"),
            ("[그림] 조합별 점수표 / 에어포일 순위표", 0, "todo"),
        ],
    )

    add_slide(
        prs,
        SEC2 + " > 2-1",
        "안정성 설계 : V-tail & 스톨 방지",
        "FC 제어에 기대지 않는 기체 고유 안정성",
        [
            ("V-tail 형상 선정", 0, "key"),
            ("웨이포인트 통과오차 최소화를 위해, FC 제어가 있음에도 동적 안정성이 보장되는", 1, None),
            ("최소 크기와 CG 모멘트암을 6개 후보군 중에서 정량 선정", 1, None),
            ("에일러론 스톨 방지", 0, "key"),
            ("elliptical 날개 + 팁에 음(-)의 incidence angle 부여", 1, None),
            ("검증 (20 m/s, AOA 10°) : 루트 부근은 뒷전 박리가 시작되지만,", 1, None),
            ("팁 부근은 유효 받음각 감소 → 국부 양력계수 감소 → 박리 지연 → 층류 유지", 2, None),
            ("→ 스톨 진입 시에도 에일러론이 살아 있어 롤 제어권을 잃지 않는다", 2, "key"),
            ("[그림] 스팬 위치별 국부 양력계수 / 주익 압력분포 / speed contour", 0, "todo"),
        ],
        note="메시지 : '제어가 잡아주니까'가 아니라, 기체 자체가 안 넘어가도록 설계했다",
    )

    add_slide(
        prs,
        SEC2 + " > 2-1",
        "구조 · 진동 해석 검증",
        "최악 하중 조건에서의 안전계수 · 공진 회피",
        [
            ("해석 조건 = 기체가 마주할 수 있는 가장 가혹한 환경", 0, "key"),
            ("최고속도 80 km/h · 최대받음각 15° · 최대추력 (쿼드 12 kgf, 푸셔 4 kgf)", 1, None),
            ("재료 모델링 : 카본파이프 → CFRP / ASA AERO → 발포·저 infill 반영해 영률·포아송비 조정", 0, None),
            ("구조해석 : fluent 공력하중 + 모터 추력 → static structural", 0, None),
            ("→ 두 재질 모두 가혹 조건에서 충분한 안전계수 확보", 1, "key"),
            ("진동해석 : 모드 20개, 블레이드 통과 주파수(BPF)로 회피 주파수 대역 산정", 0, None),
            ("유효질량비 5% 초과 & 회피대역 내 위험모드 4·5번 식별, harmonic response로 재확인", 1, None),
            ("한계도 그대로 보고", 0, None),
            ("ANSYS student 메시 제한으로 2-way FSI 대신 1-way FSI 수행", 1, None),
            ("위험모드 원인은 팁 카본파이프 지오메트리의 메싱 에러 누락 → ASA AERO가 하중 전담", 1, None),
            ("설계대로 재해석 시 공진 주파수 상승으로 회피대역 이탈 예상 (재해석 예정)", 1, None),
            ("[그림] 하중에 의한 변형 / 모터별 회피 주파수 / 위험 모드 형상", 0, "todo"),
        ],
    )

    add_slide(
        prs,
        SEC2 + " > 2-1",
        "제작 및 실물 검증",
        "설계안이 아니라 실제로 나는 기체",
        [
            ("수리1호 제작 완료 — 3D 프린팅 기반 자체 제작", 0, "key"),
            ("재료 선정 : PLA / LW-PLA / PETG / ASA / ASA-CF / ASA AERO 후보 중 ASA AERO 채택", 1, None),
            ("유리전이 95°C (LW-PLA 55°C), UV 저항·인성 우수, PLA 밀도의 절반 수준", 2, None),
            ("역할 분담 : 카본파이프 + 일반 ASA = 구조 하중 / ASA AERO = 비구조부 경량화", 1, None),
            ("25% 정렬 elliptical wing의 복잡 형상 구현에 3D 프린팅이 최적이라 판단", 1, None),
            ("완료된 실물 검증", 0, "key"),
            ("Dual Holybro M10 GPS 기반 호버링 시험 완료", 1, None),
            ("랜딩기어 retraction 동작 확인 / V-tail 조립 완료", 1, None),
            ("[사진] 호버링 / 기체 조립 / 랜딩기어 retraction / V-tail", 0, "todo"),
            ("[빈칸] 실측 제원표 — MTOW, 익폭, 전장, 공허중량, 배터리 사양", 0, "todo"),
        ],
    )

    add_slide(
        prs,
        SEC2 + " > 2-1",
        "자체 제작 시험 기체군",
        "본기 위험을 줄이기 위해 시험기까지 직접 만들었다",
        [
            ("본기(수리1호)에 곧바로 위험한 시험을 하지 않기 위해 별도 시험 기체를 자체 제작", 0, "key"),
            ("MC 테스트기체의 의미 : 최종기체와 동일한 위치 기반 제어 신호 체계를 쓰도록 제작", 0, "key"),
            ("→ 여기서 검증한 비행 로직이 코드 변경 없이 그대로 VTOL 본기로 올라간다", 1, None),
            ("→ 소프트웨어 검증 사이클을 본기 손상 위험 없이 반복 가능", 1, None),
        ],
        table=[
            ["시험 기체", "목적", "현황"],
            ["작은 수리", "수리1호 축소형 — 형상·천이 거동 저비용 검증", "[팀장 확인]"],
            ["MC 테스트기체", "자율비행 SW 검증 전용 (Pixhawk 6C + RPi5 동일 구성)", "실비행 59회분 누적"],
            ["아기수리", "조난자 구조 장치 (4장에서 상술)", "[팀장 확인]"],
        ],
        note="[빈칸] 작은 수리 제원 · 사진 · 시험 내역",
    )

    add_slide(
        prs,
        SEC2 + " > 2-2",
        "미완성 부분 개발 / 시험 계획",
        "남은 것과 그 일정",
        [],
        table=[
            ["항목", "현재 상태", "개발 / 시험 계획"],
            ["VTOL 천이–역천이 실비행", "SITL 전 구간 검증 완료, 실기체 미실시", "[빈칸] 월 [빈칸] 주 실시"],
            ["정밀착륙 폐루프", "객체인식 측 완료, 비행제어 소비부 연결 중", "연결 후 SITL → 실비행 검증"],
            ["구조장치 통합", "[팀장 확인]", "[팀장 확인]"],
            ["카메라 정밀 캘리브레이션", "근사 파라미터로 운용 중", "예선 후 실측 캘리브레이션 수행"],
            ["구조 재해석", "메싱 에러 구간 존재", "재해석 후 공진 회피 재확인"],
        ],
        note="※ 계획 열의 날짜를 반드시 채울 것 — 날짜 없는 계획은 계획으로 읽히지 않는다",
    )

    add_slide(
        prs,
        SEC2 + " > 2-3",
        "안전장치 구현 상황",
        "물리 · 비행제어 · 소프트웨어 3중 안전망",
        [
            ("① 물리 안전장치 — CO2 카트리지 격발 낙하산", 0, "key"),
            ("화약식 대비 즉각적인 기계적 사출로 전개 시점을 극적으로 앞당김", 1, None),
            ("저고도 자유낙하에서도 신속히 유효 항력 형성 → 충격 에너지 격감", 1, None),
            ("② 비행제어 안전 알고리즘", 0, "key"),
            ("복귀 가능 에너지 동적 산출 : 단순 잔량(%) 페일세이프에서 탈피", 1, None),
            ("MC 기동전력 1000W + FW 순항(70km/h, 전비 400m/Wh) 기반 실시간 귀환 요구 에너지 연산", 2, None),
            ("6S 10000mAh의 20% 안전마진 도달 즉시 자율 귀환(RTL)", 2, None),
            ("비상 활공 천이 : MC 구동계 고장으로 자세제어 불가 시, 자세를 억지로 잡는 대신", 1, None),
            ("의도적 급강하로 위치E→운동E 전환 → 실속속도 확보 → FW 천이 → 푸셔 최대추력 → 귀환·활공착륙", 2, None),
            ("③ 소프트웨어 안전망 (실비행에서 실제 발동한 것들)", 0, "key"),
            ("RC 모드스위치 + ROS 토픽 이중 오버라이드 → 실패 시 AUTO.LOITER 폴백", 1, None),
            ("조종사 인계 감지 시 자율노드 자동 정지 (실비행 사고 대응으로 신설)", 1, None),
            ("상태별 탈출 타임아웃 4종 + 이륙지점 기준 거리상한 300m (마지막 기하학적 제동장치)", 1, None),
            ("[그림] 페일세이프 다이어그램 / 글라이딩 시뮬레이션 결과", 0, "todo"),
        ],
    )

    # ---------------- 3. 자동비행 시스템 ----------------
    add_slide(
        prs,
        SEC3 + " > 3-1",
        "유도 · 제어 · 항법 아키텍처",
        "3계층 분리 · 알고리즘 교체 가능",
        [
            ("계층 구조 : 경로생성(Path Planner) → 경로추종(Path Controller) → 자세제어(Inner Loop, FC)", 0, "key"),
            ("1차의 '완전한 모듈화 · 자율적 알고리즘 대체' 설계가 실제 코드 구조로 실현됨", 0, None),
            ("ROS2 노드 3종 (Telemetry / Offboard / Mission) + ROS 비의존 순수 라이브러리로 분리", 1, None),
            ("계층마다 구현체를 교체해도 나머지 계층이 바뀌지 않는다", 1, None),
            ("소프트웨어 규모 : 3개 도메인 약 51,400줄 / 205개 파일 (비행제어 · 객체인식 · 시뮬레이터)", 0, None),
            ("[그림] block diagram of navigation process (1차 fig.2.1.1 — 캡션을 '구현 결과'로 수정)", 0, "todo"),
        ],
    )

    add_slide(
        prs,
        SEC3 + " > 3-1",
        "경로생성 : eta³ Clothoid",
        "11종을 직접 만들어 비교하고 골랐다",
        [
            ("곡률 연속(G2) 경로 — clothoid는 곡률이 호길이에 선형 증가하여 조종면 급변이 없다", 0, "key"),
            ("자체 최적화 방법론", 0, None),
            ("양단 G2 보간을 따로 하지 않고 NR 잔차와 동시 계산 → 위치 잔차와 매개변수를 종속화", 1, None),
            ("위치·헤딩 잔차에 가중치를 부여해 전 구간에 분산 → 종단 WP 이탈 문제 해결", 1, None),
            ("곡률제한 실효성이 없는 구간에는 가상 WP를 삽입해 제한 문제 자체를 회피", 1, None),
            ("BasePlanner 구현체 11종(clothoid, D-iterpin 등)을 제작·시각화 비교하여 최우수안 선정", 0, "key"),
            ("파라미터도 실측으로 확정 : v_cruise 18.0 = 실측 FW_AIRSPD_TRIM 15 + 배풍 마진 3", 0, None),
            ("경로 생성 비용·품질이 특정 구간에서 계단형으로 갈린다는 것을 실측으로 확인", 1, None),
            ("[그림] 경로 생성 6패널 플롯 / planner 11종 비교 시각화", 0, "todo"),
        ],
    )

    add_slide(
        prs,
        SEC3 + " > 3-1",
        "경로추종 및 비행 상태머신",
        "이륙부터 착륙까지 12상태 자동 전이",
        [
            ("L1 유도 기반 경로추종 + 위치 setpoint 채널", 0, "key"),
            ("고정익 오프보드는 속도 setpoint를 무시한다는 사실을 SITL 실측으로 확정한 결과", 1, None),
            ("상태 전이 : ARM_TAKEOFF → CLIMBING → TRANSITION_FW → STREAMING → FOLLOWING", 0, None),
            ("→ TRANSITION_MC → HOLD → LANDING → DONE  (+ OVERRIDE / PILOT_TAKEOVER)", 1, None),
            ("vehicle_type 파라미터 하나로 MC 테스트기체와 VTOL 본기가 동일 코드를 사용", 0, "key"),
            ("→ 2장의 시험기체 전략이 코드 수준에서 실제로 성립하는 근거", 1, None),
            ("[그림] 상태 전이도", 0, "todo"),
        ],
    )

    add_slide(
        prs,
        SEC3 + " > 3-2",
        "검증 체계 : 3단 검증 피라미드",
        "검증을 통과하지 못한 코드는 기체에 올라가지 않는다",
        [
            ("[3층] 실기체 비행 — 59회분 비행 로그 · ulog 79건 전량 자동 분석", 0, "key"),
            ("[2층] SITL 회귀 — 시나리오 24런 자동 실행 · 자동 판정", 0, "key"),
            ("[1층] 단위 테스트 — 1,811건 (비행제어 783 + 객체인식 1,028), ROS 없이 실행", 0, "key"),
            ("구조가 검증을 가능하게 했다 — 비행 로직을 ROS 비의존 순수함수로 분리했기에 1층이 성립", 0, None),
            ("SITL 성과 사례 : PX4 상류 회귀 버그를 소스 계보까지 추적해 2줄 패치로 해소", 0, None),
            ("북향 이탈 21.78 m → 0.38 m / 고도편차 판정 FAIL → PASS / 회귀 0건", 1, "key"),
            ("비행 로그는 매 비행 전량 회수 → 자동 진단서 생성 → 결함 규명 → 수정 → 재배포", 0, None),
            ("[그림] 3단 피라미드 도식 / SITL 자동 판정서 캡처", 0, "todo"),
        ],
        note="발표 시간 배분을 가장 크게 줄 슬라이드 — 질의응답 방어력이 여기서 나온다",
    )

    add_slide(
        prs,
        SEC3 + " > 3-2",
        "지상 및 비행 시험 현황 / 계획",
        "누적 59회분 + 남은 시험",
        [
            ("완료 : 호버링 시험(Dual M10 GPS) / MC 자율비행 59회분 / SITL 전 사이클 PASS", 0, "key"),
            ("대표 결과 : 조종사 개입 0으로 완주 — 경로이탈 최대 0.94 m, WP 최근접 0.10~0.64 m", 0, "key"),
            ("개발 사이클 : 비행 → 로그 전량 회수 → 자동 분석 → 결함 규명 → 수정 → 재배포", 0, None),
            ("남은 시험", 0, None),
            ("VTOL 천이–역천이 실비행 / 정밀착륙 통합 / 구조장치 연동 / 통합 리허설", 1, None),
            ("[그림] 비행 궤적 플롯 / 자동 생성 분석서 캡처", 0, "todo"),
        ],
    )

    # ---------------- 4. 임무수행 구현기술 ----------------
    add_slide(
        prs,
        SEC4 + " > 4-1",
        "임무수행 시나리오",
        "이륙에서 구조 · 복귀까지",
        [
            ("이륙 → MC→FW 천이 → 경로비행 → 목표구역 도착 → FW→MC 역천이", 0, "key"),
            ("→ 구역 스캔 → 타겟 포착·확인 → 정중앙 상공 정렬 → 정밀착륙", 0, "key"),
            ("→ 아기수리 전개 → 조난자 인양 → 모선 복귀 → 이륙 → 복귀·착륙", 0, "key"),
        ],
        table=[
            ["임무 타겟", "규격", "착륙 지점"],
            ["버티포트", "흰 필드 Ø3 m + 빨간 원 Ø2 m + 중앙 ArUco 50 cm", "정중앙"],
            ["조난자 구역", "초록 매트 3.0 × 3.0 × 0.105 m + 중앙 흰 박스 20 cm", "박스 옆 빈 초록면"],
            ["하기구역", "빨간 십자 (대회측 규격 비공개)", "중앙"],
            ["단순 착륙", "피듀셜 없음", "GPS + 라이다 하강"],
        ],
    )

    add_slide(
        prs,
        SEC4 + " > 4-2",
        "임무 장치 : 아기수리",
        "모선과 분리된 전용 구조선",
        [
            ("설계 사상 : 모선이 직접 인양하는 방식은 정확도·신뢰도가 떨어진다", 0, "key"),
            ("→ 구조 전용 기체를 별도로 설계·제작", 1, None),
            ("동작 순서 : 수리1호가 조난자 근처 착륙 → 아기수리 전개", 0, None),
            ("지면 주행 (서보모터 4개 조향) → 트레이 인양 (양측 그리퍼) → 모선 복귀 (상단 그리퍼)", 1, None),
            ("[팀장 담당 · 전면 빈칸]", 0, "todo"),
            ("제작 현황 사진 / 구동 시험 결과 / 1차 설계 대비 변경점 / 모선 탑재·분리 메커니즘", 1, "todo"),
            ("[그림] 아기수리 구조 과정 다이어그램(1차) + 실물 사진", 0, "todo"),
        ],
    )

    add_slide(
        prs,
        SEC4 + " > 4-2",
        "객체인식 및 정밀착륙",
        "결정론적 고전 CV + 상대 pose 폐루프",
        [
            ("고전 CV(OpenCV) 채택 근거 : 결정론 = 신뢰도 / RPi5 경량 / 학습 데이터셋 부담 0", 0, "key"),
            ("→ 대회 당일 '오작동·오인식 절대 불가' 요구에 정면 대응", 1, None),
            ("coarse → fine 2단 검출 : 고고도에서 소형 피처는 원리적으로 디코드 불가 (GSD 분석)", 0, "key"),
            ("→ fine 피처 락·검증 전에는 최종 착륙을 커밋하지 않는다 (코드 구조로 강제)", 1, None),
            ("30 cm 정밀도는 폐루프로만 달성 — GPS 절대좌표는 1~2 m가 하한", 0, "key"),
            ("카메라가 제어오차를 직접 측정하는 비주얼 서보잉 방식", 1, None),
            ("실측 성능 : 프레임당 처리 18 ms, solvePnP 거리오차 0.25~0.58 %", 0, None),
            ("착륙점은 기체 최외곽 반경 기반 안전창에서 산출 (물리 치수 반영, 매트 이탈 차단)", 0, None),
            ("[그림] 검출 오버레이 / 고도별 GSD 표 / 데이터 흐름도", 0, "todo"),
        ],
    )

    add_slide(
        prs,
        SEC4 + " > 4-3",
        "보완 계획",
        "본선까지의 마일스톤",
        [
            ("임무 수행 측 보완 항목 (시스템·기체 측은 2-2 참조)", 0, None),
        ],
        table=[
            ["마일스톤", "내용", "목표 시기"],
            ["구조장치 통합 시험", "아기수리 탑재·전개·인양 통합 검증", "[빈칸]"],
            ["정밀착륙 폐루프 실비행", "비행제어 소비부 연결 후 실기체 검증", "[빈칸]"],
            ["검출기 실촬영 재튜닝", "실제 자갈·조명 조건 데이터로 검출 임계 재조정", "[빈칸]"],
            ["구조 재해석", "메싱 에러 해소 후 공진 회피 재확인", "[빈칸]"],
            ["통합 리허설", "임무 전 구간 연속 수행", "[빈칸]"],
        ],
        note="※ 목표 시기를 반드시 채울 것",
    )

    # ---------------- 5. 기타 ----------------
    add_slide(
        prs,
        SEC5 + " > 5-1",
        "비행시험 데이터 & 영상",
        "제출 자료 안내",
        [
            ("제출 비행데이터 — 각 데이터가 무엇을 보여주는지 1줄씩", 0, None),
            ("[빈칸] 제출 대상 로그 선정", 1, "todo"),
            ("제출 영상 — 전체시스템 / 구성별 / 안전장치 / 지상국", 0, None),
            ("[빈칸] 촬영 및 편집", 1, "todo"),
            ("", 0, None),
            ("맺음말", 0, "key"),
            ("1. 기체 · 소프트웨어 · 시험장비까지 전부 자체 제작", 1, None),
            ("2. 해석으로 설계하고, 3단 검증 체계로 신뢰성을 확보", 1, None),
            ("3. 모든 구성요소가 교체 가능한 구조 — 본선까지 계속 개선 가능", 1, None),
        ],
    )

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    print(f"저장: {path}  (슬라이드 {len(prs.slides.__iter__.__self__._sldIdLst)}장)")


if __name__ == "__main__":
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("docs/qualifier2_report_draft.pptx")
    build(out)
